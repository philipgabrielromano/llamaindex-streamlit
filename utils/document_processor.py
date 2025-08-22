# utils/document_processor.py (Fixed without LlamaIndex text splitter)
import streamlit as st
import hashlib
import mimetypes
from datetime import datetime
from typing import List, Dict, Optional, Tuple
import PyPDF2
import docx
import pandas as pd
import io
from pathlib import Path
import json
import xml.etree.ElementTree as ET
import re

# Try to import additional libraries
try:
    import pptx
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

# Simple Document class if LlamaIndex isn't available
class SimpleDocument:
    """Simple document class as fallback"""
    def __init__(self, text: str, metadata: Dict = None):
        self.text = text
        self.metadata = metadata or {}

# Try to import LlamaIndex Document, fallback to simple version
try:
    from llama_index.core import Document
except ImportError:
    try:
        from llama_index import Document
    except ImportError:
        Document = SimpleDocument

class SimpleTextSplitter:
    """Simple text splitter implementation without LlamaIndex dependency"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200, separators: List[str] = None):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = separators or ["\n\n", "\n", ". ", " ", ""]
    
    def split_text(self, text: str) -> List[str]:
        """Split text into chunks using the specified separators"""
        if not text or len(text) <= self.chunk_size:
            return [text] if text else []
        
        chunks = []
        current_chunk = ""
        
        # Split by separators in order of preference
        splits = [text]
        
        for separator in self.separators:
            new_splits = []
            for split in splits:
                if len(split) > self.chunk_size:
                    new_splits.extend(split.split(separator))
                else:
                    new_splits.append(split)
            splits = new_splits
        
        # Combine splits into chunks
        for split in splits:
            if len(current_chunk) + len(split) <= self.chunk_size:
                current_chunk += split
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                
                # Handle overlap
                if self.chunk_overlap > 0 and chunks:
                    overlap_text = chunks[-1][-self.chunk_overlap:] if len(chunks[-1]) > self.chunk_overlap else chunks[-1]
                    current_chunk = overlap_text + split
                else:
                    current_chunk = split
        
        # Add the last chunk
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        # Filter out empty chunks
        return [chunk for chunk in chunks if chunk.strip()]

class DocumentProcessor:
    """Enhanced document processor without LlamaIndex text splitter dependency"""
    
    # Define supported file types and their processors
    SUPPORTED_TYPES = {
        'pdf': 'PDF Documents',
        'docx': 'Word Documents', 
        'txt': 'Text Files',
        'md': 'Markdown Files',
        'html': 'HTML Files',
        'csv': 'CSV Files',
        'json': 'JSON Files',
        'xml': 'XML Files',
        'pptx': 'PowerPoint Presentations',
        'xlsx': 'Excel Spreadsheets',
        'rtf': 'Rich Text Format'
    }
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = SimpleTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
    
    def get_supported_types(self) -> Dict[str, str]:
        """Get dictionary of supported file types"""
        return self.SUPPORTED_TYPES.copy()
    
    def is_supported_type(self, filename: str) -> bool:
        """Check if file type is supported"""
        file_ext = self._get_file_type(filename)
        return file_ext in self.SUPPORTED_TYPES
    
    def process_uploaded_file(self, uploaded_file) -> Optional[Document]:
        """Process a single uploaded file and return a Document object"""
        try:
            content = uploaded_file.read()
            file_type = self._get_file_type(uploaded_file.name)
            
            if not self.is_supported_type(uploaded_file.name):
                st.error(f"Unsupported file type: {file_type}")
                return None
            
            # Extract text based on file type
            text = self._extract_text_by_type(content, file_type, uploaded_file.name)
            
            if not text or not text.strip():
                st.warning(f"No text content extracted from {uploaded_file.name}")
                return None
            
            # Create document with metadata
            document = Document(
                text=text,
                metadata={
                    'filename': uploaded_file.name,
                    'file_type': file_type,
                    'file_size': len(content),
                    'processed_at': datetime.now().isoformat(),
                    'chunk_size': self.chunk_size,
                    'chunk_overlap': self.chunk_overlap,
                    'document_hash': self._generate_hash(content),
                    'source': 'manual_upload',
                    'text_length': len(text),
                    'word_count': len(text.split())
                }
            )
            
            return document
            
        except Exception as e:
            st.error(f"Error processing file {uploaded_file.name}: {str(e)}")
            return None
    
    def process_sharepoint_documents(self, documents: List[Dict]) -> List[Document]:
        """Process documents from SharePoint"""
        processed_docs = []
        
        for doc_info in documents:
            try:
                content = doc_info.get('content', '')
                if not content or not content.strip():
                    st.warning(f"No content found in {doc_info.get('filename', 'Unknown file')}")
                    continue
                
                # Create document from SharePoint info
                document = Document(
                    text=content,
                    metadata={
                        'filename': doc_info.get('filename', 'Unknown'),
                        'sharepoint_id': doc_info.get('id'),
                        'modified_date': doc_info.get('modified'),
                        'file_path': doc_info.get('file_path'),
                        'processed_at': datetime.now().isoformat(),
                        'chunk_size': self.chunk_size,
                        'chunk_overlap': self.chunk_overlap,
                        'source': 'sharepoint',
                        'text_length': len(content),
                        'word_count': len(content.split())
                    }
                )
                processed_docs.append(document)
                
            except Exception as e:
                st.error(f"Error processing SharePoint document {doc_info.get('filename', 'Unknown')}: {str(e)}")
        
        return processed_docs
    
    def _extract_text_by_type(self, content: bytes, file_type: str, filename: str) -> str:
        """Extract text based on file type"""
        try:
            if file_type == 'pdf':
                return self._extract_pdf_text(content)
            elif file_type == 'docx':
                return self._extract_docx_text(content)
            elif file_type == 'txt':
                return self._extract_txt_text(content)
            elif file_type == 'md':
                return self._extract_markdown_text(content)
            elif file_type == 'html':
                return self._extract_html_text(content)
            elif file_type == 'csv':
                return self._extract_csv_text(content)
            elif file_type == 'json':
                return self._extract_json_text(content)
            elif file_type == 'xml':
                return self._extract_xml_text(content)
            elif file_type == 'pptx' and PPTX_AVAILABLE:
                return self._extract_pptx_text(content)
            elif file_type == 'xlsx':
                return self._extract_xlsx_text(content)
            else:
                # Fallback: try to decode as text
                return content.decode('utf-8', errors='ignore')
                
        except Exception as e:
            raise Exception(f"Text extraction failed for {file_type}: {str(e)}")
    
    def _extract_pdf_text(self, content: bytes) -> str:
        """Extract text from PDF content"""
        try:
            pdf_file = io.BytesIO(content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text = ""
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                except Exception as e:
                    st.warning(f"Could not extract text from page {page_num + 1}")
            
            return text.strip()
            
        except Exception as e:
            raise Exception(f"PDF extraction failed: {str(e)}")
    
    def _extract_docx_text(self, content: bytes) -> str:
        """Extract text from DOCX content"""
        try:
            docx_file = io.BytesIO(content)
            doc = docx.Document(docx_file)
            
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            return text.strip()
            
        except Exception as e:
            raise Exception(f"DOCX extraction failed: {str(e)}")
    
    def _extract_txt_text(self, content: bytes) -> str:
        """Extract text from plain text files"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'utf-16', 'latin-1', 'cp1252']:
                try:
                    return content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            
            # Fallback with error handling
            return content.decode('utf-8', errors='ignore')
            
        except Exception as e:
            raise Exception(f"Text extraction failed: {str(e)}")
    
    def _extract_markdown_text(self, content: bytes) -> str:
        """Extract text from Markdown files"""
        try:
            text = content.decode('utf-8', errors='ignore')
            
            # Remove markdown syntax (basic cleanup)
            # Remove headers
            text = re.sub(r'^#+\s+', '', text, flags=re.MULTILINE)
            # Remove bold/italic
            text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
            text = re.sub(r'\*([^*]+)\*', r'\1', text)
            # Remove links but keep text
            text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
            # Remove code blocks
            text = re.sub(r'```[^`]*```', '', text, flags=re.DOTALL)
            text = re.sub(r'`([^`]+)`', r'\1', text)
            
            return text.strip()
            
        except Exception as e:
            raise Exception(f"Markdown extraction failed: {str(e)}")
    
    def _extract_html_text(self, content: bytes) -> str:
        """Extract text from HTML content"""
        try:
            html_content = content.decode('utf-8', errors='ignore')
            
            if BS4_AVAILABLE:
                soup = BeautifulSoup(html_content, 'html.parser')
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text(separator='\n', strip=True)
            else:
                # Basic HTML tag removal
                text = re.sub(r'<[^>]+>', '', html_content)
                text = re.sub(r'&[a-zA-Z0-9#]+;', ' ', text)  # Remove HTML entities
                return text.strip()
                
        except Exception as e:
            raise Exception(f"HTML extraction failed: {str(e)}")
    
    def _extract_csv_text(self, content: bytes) -> str:
        """Extract text from CSV content"""
        try:
            csv_text = content.decode('utf-8', errors='ignore')
            df = pd.read_csv(io.StringIO(csv_text))
            
            # Convert dataframe to text representation
            text = df.to_string(index=False)
            return text
            
        except Exception as e:
            raise Exception(f"CSV extraction failed: {str(e)}")
    
    def _extract_json_text(self, content: bytes) -> str:
        """Extract text from JSON content"""
        try:
            json_text = content.decode('utf-8', errors='ignore')
            data = json.loads(json_text)
            
            # Convert JSON to readable text
            def extract_values(obj, prefix=""):
                text_parts = []
                if isinstance(obj, dict):
                    for key, value in obj.items():
                        if isinstance(value, (dict, list)):
                            text_parts.extend(extract_values(value, f"{prefix}{key}: "))
                        else:
                            text_parts.append(f"{prefix}{key}: {str(value)}")
                elif isinstance(obj, list):
                    for i, item in enumerate(obj):
                        if isinstance(item, (dict, list)):
                            text_parts.extend(extract_values(item, f"{prefix}[{i}] "))
                        else:
                            text_parts.append(f"{prefix}[{i}]: {str(item)}")
                else:
                    text_parts.append(f"{prefix}{str(obj)}")
                return text_parts
            
            text_parts = extract_values(data)
            return "\n".join(text_parts)
            
        except Exception as e:
            raise Exception(f"JSON extraction failed: {str(e)}")
    
    def _extract_xml_text(self, content: bytes) -> str:
        """Extract text from XML content"""
        try:
            xml_text = content.decode('utf-8', errors='ignore')
            root = ET.fromstring(xml_text)
            
            # Extract all text content from XML
            def extract_text(element):
                text_parts = []
                if element.text:
                    text_parts.append(element.text.strip())
                for child in element:
                    text_parts.extend(extract_text(child))
                if element.tail:
                    text_parts.append(element.tail.strip())
                return text_parts
            
            text_parts = extract_text(root)
            return "\n".join([part for part in text_parts if part])
            
        except Exception as e:
            raise Exception(f"XML extraction failed: {str(e)}")
    
    def _extract_pptx_text(self, content: bytes) -> str:
        """Extract text from PPTX content"""
        try:
            pptx_file = io.BytesIO(content)
            presentation = pptx.Presentation(pptx_file)
            
            text = ""
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = f"Slide {slide_num + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                
                text += slide_text + "\n"
            
            return text.strip()
            
        except Exception as e:
            raise Exception(f"PPTX extraction failed: {str(e)}")
    
    def _extract_xlsx_text(self, content: bytes) -> str:
        """Extract text from XLSX content"""
        try:
            xlsx_file = io.BytesIO(content)
            
            # Read all sheets
            excel_data = pd.read_excel(xlsx_file, sheet_name=None)
            
            text = ""
            for sheet_name, df in excel_data.items():
                text += f"Sheet: {sheet_name}\n"
                text += df.to_string(index=False) + "\n\n"
            
            return text.strip()
            
        except Exception as e:
            raise Exception(f"XLSX extraction failed: {str(e)}")
    
    def _get_file_type(self, filename: str) -> str:
        """Get file type from filename"""
        return Path(filename).suffix.lower().lstrip('.')
    
    def _generate_hash(self, content: bytes) -> str:
        """Generate MD5 hash of file content"""
        return hashlib.md5(content).hexdigest()
    
    def chunk_text(self, text: str) -> List[str]:
        """Manually chunk text if needed"""
        return self.text_splitter.split_text(text)
    
    def get_document_stats(self, document) -> Dict:
        """Get statistics about a document"""
        if hasattr(document, 'text'):
            text = document.text
        elif isinstance(document, dict):
            text = document.get('content', '')
        else:
            text = str(document)
        
        text_length = len(text)
        word_count = len(text.split())
        estimated_chunks = max(1, text_length // self.chunk_size)
        
        filename = 'Unknown'
        if hasattr(document, 'metadata'):
            filename = document.metadata.get('filename', 'Unknown')
        elif isinstance(document, dict):
            filename = document.get('filename', 'Unknown')
        
        return {
            'text_length': text_length,
            'word_count': word_count,
            'estimated_chunks': estimated_chunks,
            'filename': filename
        }
    
    def create_document_from_text(self, text: str, metadata: Dict = None) -> Document:
        """Create a document object from text and metadata"""
        return Document(
            text=text,
            metadata=metadata or {}
        )
    
    def update_chunk_settings(self, chunk_size: int, chunk_overlap: int):
        """Update chunking settings"""
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = SimpleTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )
